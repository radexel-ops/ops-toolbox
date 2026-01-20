"""File Upload API Router"""

from fastapi import APIRouter, UploadFile, File, HTTPException
from typing import List, Dict, Any
import base64
import os
import uuid
import io
from pathlib import Path

# PDF 텍스트 추출
try:
    from PyPDF2 import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("[Warning] PyPDF2 not installed, PDF text extraction disabled")

# Word (DOCX) 텍스트 추출
try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("[Warning] python-docx not installed, DOCX text extraction disabled")

# Excel (XLSX) 텍스트 추출
try:
    from openpyxl import load_workbook
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    print("[Warning] openpyxl not installed, XLSX text extraction disabled")

# HWP 텍스트 추출 (한글 워드프로세서)
try:
    import olefile
    HWP_SUPPORT = True
except ImportError:
    HWP_SUPPORT = False
    print("[Warning] olefile not installed, HWP text extraction disabled")

# PowerPoint (PPTX) 텍스트 추출
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("[Warning] python-pptx not installed, PPTX text extraction disabled")

# 구형 Excel (XLS) 텍스트 추출
try:
    import xlrd
    XLS_SUPPORT = True
except ImportError:
    XLS_SUPPORT = False
    print("[Warning] xlrd not installed, XLS text extraction disabled")

# OpenDocument (ODT, ODS, ODP) 텍스트 추출
try:
    from odf import text as odf_text
    from odf.opendocument import load as odf_load
    ODF_SUPPORT = True
except ImportError:
    ODF_SUPPORT = False
    print("[Warning] odfpy not installed, ODF text extraction disabled")

# RTF 텍스트 추출
try:
    from striprtf.striprtf import rtf_to_text
    RTF_SUPPORT = True
except ImportError:
    RTF_SUPPORT = False
    print("[Warning] striprtf not installed, RTF text extraction disabled")

router = APIRouter()


def extract_pdf_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """PDF 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not PDF_SUPPORT:
        return "", False

    try:
        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)
        text_parts = []
        total_chars = 0
        truncated = False

        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if total_chars + len(page_text) > max_chars:
                remaining = max_chars - total_chars
                text_parts.append(f"--- 페이지 {page_num + 1} ---\n{page_text[:remaining]}")
                truncated = True
                break
            text_parts.append(f"--- 페이지 {page_num + 1} ---\n{page_text}")
            total_chars += len(page_text)

        return "\n\n".join(text_parts), truncated
    except Exception as e:
        print(f"[PDF Extract Error] {e}")
        return f"[PDF 텍스트 추출 실패: {str(e)}]", False


def extract_docx_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """DOCX 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not DOCX_SUPPORT:
        return "", False

    try:
        docx_file = io.BytesIO(content)
        doc = DocxDocument(docx_file)
        text_parts = []
        total_chars = 0
        truncated = False

        for para in doc.paragraphs:
            para_text = para.text
            if total_chars + len(para_text) > max_chars:
                remaining = max_chars - total_chars
                text_parts.append(para_text[:remaining])
                truncated = True
                break
            text_parts.append(para_text)
            total_chars += len(para_text) + 1  # +1 for newline

        # 테이블 내용도 추출
        if not truncated:
            for table in doc.tables:
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        row_line = " | ".join(row_texts)
                        if total_chars + len(row_line) > max_chars:
                            truncated = True
                            break
                        text_parts.append(row_line)
                        total_chars += len(row_line) + 1
                if truncated:
                    break

        return "\n".join(text_parts), truncated
    except Exception as e:
        print(f"[DOCX Extract Error] {e}")
        return f"[DOCX 텍스트 추출 실패: {str(e)}]", False


def extract_xlsx_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """XLSX 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not XLSX_SUPPORT:
        return "", False

    try:
        xlsx_file = io.BytesIO(content)
        wb = load_workbook(xlsx_file, read_only=True, data_only=True)
        text_parts = []
        total_chars = 0
        truncated = False

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"--- 시트: {sheet_name} ---")
            total_chars += len(sheet_name) + 12

            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if total_chars > max_chars:
                    truncated = True
                    break

                row_texts = [str(cell) if cell is not None else "" for cell in row]
                row_line = " | ".join(row_texts)

                if total_chars + len(row_line) > max_chars:
                    truncated = True
                    break

                text_parts.append(row_line)
                total_chars += len(row_line) + 1
                row_count += 1

                # 시트당 최대 1000행
                if row_count >= 1000:
                    text_parts.append(f"... (최대 1000행 표시)")
                    break

            if truncated:
                break

        wb.close()
        return "\n".join(text_parts), truncated
    except Exception as e:
        print(f"[XLSX Extract Error] {e}")
        return f"[XLSX 텍스트 추출 실패: {str(e)}]", False


def extract_hwp_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """HWP 파일에서 텍스트 추출 (한글 워드프로세서). 반환: (텍스트, 잘림여부)"""
    if not HWP_SUPPORT:
        return "", False

    try:
        hwp_file = io.BytesIO(content)
        ole = olefile.OleFileIO(hwp_file)

        # HWP 파일 구조에서 텍스트 추출
        # 'PrvText' 또는 'BodyText/Section#' 스트림에서 텍스트 추출
        text_parts = []
        total_chars = 0
        truncated = False

        # 미리보기 텍스트 (빠른 추출)
        if ole.exists('PrvText'):
            prv_data = ole.openstream('PrvText').read()
            try:
                # UTF-16 LE 디코딩
                text = prv_data.decode('utf-16-le', errors='ignore')
                text = text.replace('\x00', '').strip()
                if text:
                    if len(text) > max_chars:
                        text = text[:max_chars]
                        truncated = True
                    text_parts.append(text)
                    total_chars += len(text)
            except:
                pass

        # BodyText 섹션에서 추가 텍스트 추출
        if not truncated:
            section_num = 0
            while ole.exists(f'BodyText/Section{section_num}'):
                section_data = ole.openstream(f'BodyText/Section{section_num}').read()
                try:
                    # 압축 해제 시도 (zlib)
                    import zlib
                    try:
                        decompressed = zlib.decompress(section_data, -15)
                        section_data = decompressed
                    except:
                        pass

                    # 텍스트 추출 (간단한 방법)
                    # HWP 포맷의 텍스트는 특정 바이트 패턴 사이에 있음
                    text = section_data.decode('utf-16-le', errors='ignore')
                    text = ''.join(c for c in text if c.isprintable() or c in '\n\r\t')
                    text = text.strip()

                    if text and total_chars + len(text) <= max_chars:
                        text_parts.append(text)
                        total_chars += len(text)
                    elif text:
                        remaining = max_chars - total_chars
                        if remaining > 0:
                            text_parts.append(text[:remaining])
                        truncated = True
                        break
                except:
                    pass
                section_num += 1

        ole.close()

        result_text = "\n".join(text_parts)
        # 중복 제거 및 정리
        result_text = '\n'.join(line for line in result_text.split('\n') if line.strip())

        return result_text, truncated
    except Exception as e:
        print(f"[HWP Extract Error] {e}")
        return f"[HWP 텍스트 추출 실패: {str(e)}]", False


def extract_pptx_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """PPTX 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not PPTX_SUPPORT:
        return "[python-pptx가 설치되지 않아 PPTX를 읽을 수 없습니다.]", False

    try:
        pptx_file = io.BytesIO(content)
        prs = Presentation(pptx_file)
        text_parts = []
        total_chars = 0
        truncated = False

        for slide_num, slide in enumerate(prs.slides):
            slide_texts = [f"--- 슬라이드 {slide_num + 1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)

            slide_text = "\n".join(slide_texts)
            if total_chars + len(slide_text) > max_chars:
                remaining = max_chars - total_chars
                text_parts.append(slide_text[:remaining])
                truncated = True
                break
            text_parts.append(slide_text)
            total_chars += len(slide_text)

        return "\n\n".join(text_parts), truncated
    except Exception as e:
        print(f"[PPTX Extract Error] {e}")
        return f"[PPTX 텍스트 추출 실패: {str(e)}]", False


def extract_xls_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """XLS (구형 Excel) 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not XLS_SUPPORT:
        return "[xlrd가 설치되지 않아 XLS를 읽을 수 없습니다.]", False

    try:
        xls_file = io.BytesIO(content)
        wb = xlrd.open_workbook(file_contents=content)
        text_parts = []
        total_chars = 0
        truncated = False

        for sheet_name in wb.sheet_names():
            sheet = wb.sheet_by_name(sheet_name)
            text_parts.append(f"--- 시트: {sheet_name} ---")
            total_chars += len(sheet_name) + 12

            for row_idx in range(min(sheet.nrows, 1000)):  # 최대 1000행
                if total_chars > max_chars:
                    truncated = True
                    break

                row_texts = [str(cell.value) if cell.value else "" for cell in sheet.row(row_idx)]
                row_line = " | ".join(row_texts)

                if total_chars + len(row_line) > max_chars:
                    truncated = True
                    break

                text_parts.append(row_line)
                total_chars += len(row_line) + 1

            if truncated:
                break

        return "\n".join(text_parts), truncated
    except Exception as e:
        print(f"[XLS Extract Error] {e}")
        return f"[XLS 텍스트 추출 실패: {str(e)}]", False


def extract_odt_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """ODT (OpenDocument Text) 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not ODF_SUPPORT:
        return "[odfpy가 설치되지 않아 ODT를 읽을 수 없습니다.]", False

    try:
        odt_file = io.BytesIO(content)
        doc = odf_load(odt_file)
        text_parts = []
        total_chars = 0
        truncated = False

        for para in doc.getElementsByType(odf_text.P):
            para_text = ""
            for node in para.childNodes:
                if node.nodeType == node.TEXT_NODE:
                    para_text += node.data
                elif hasattr(node, 'data'):
                    para_text += node.data

            if total_chars + len(para_text) > max_chars:
                remaining = max_chars - total_chars
                text_parts.append(para_text[:remaining])
                truncated = True
                break
            text_parts.append(para_text)
            total_chars += len(para_text) + 1

        return "\n".join(text_parts), truncated
    except Exception as e:
        print(f"[ODT Extract Error] {e}")
        return f"[ODT 텍스트 추출 실패: {str(e)}]", False


def extract_rtf_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """RTF 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    if not RTF_SUPPORT:
        return "[striprtf가 설치되지 않아 RTF를 읽을 수 없습니다.]", False

    try:
        # RTF는 보통 CP1252 또는 UTF-8
        try:
            rtf_content = content.decode('utf-8')
        except:
            rtf_content = content.decode('cp1252', errors='ignore')

        text = rtf_to_text(rtf_content)
        truncated = len(text) > max_chars
        if truncated:
            text = text[:max_chars]

        return text, truncated
    except Exception as e:
        print(f"[RTF Extract Error] {e}")
        return f"[RTF 텍스트 추출 실패: {str(e)}]", False


def extract_hwpx_text(content: bytes, max_chars: int = 50000) -> tuple[str, bool]:
    """HWPX (한글 XML) 파일에서 텍스트 추출. 반환: (텍스트, 잘림여부)"""
    import zipfile
    import xml.etree.ElementTree as ET

    try:
        hwpx_file = io.BytesIO(content)
        text_parts = []
        total_chars = 0
        truncated = False

        with zipfile.ZipFile(hwpx_file, 'r') as zf:
            # HWPX는 ZIP 안에 XML 파일들이 있음
            for name in zf.namelist():
                if 'section' in name.lower() and name.endswith('.xml'):
                    xml_content = zf.read(name)
                    try:
                        root = ET.fromstring(xml_content)
                        # 모든 텍스트 노드 추출
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                text = elem.text.strip()
                                if total_chars + len(text) > max_chars:
                                    remaining = max_chars - total_chars
                                    text_parts.append(text[:remaining])
                                    truncated = True
                                    break
                                text_parts.append(text)
                                total_chars += len(text) + 1
                        if truncated:
                            break
                    except ET.ParseError:
                        continue

        result = "\n".join(text_parts)
        return result if result else "[HWPX 텍스트를 찾을 수 없습니다.]", truncated
    except Exception as e:
        print(f"[HWPX Extract Error] {e}")
        return f"[HWPX 텍스트 추출 실패: {str(e)}]", False


def extract_document_text(content: bytes, ext: str, max_chars: int = 50000) -> tuple[str, bool]:
    """문서 파일에서 텍스트 추출 (확장자별 분기). 반환: (텍스트, 잘림여부)"""
    ext = ext.lower()

    if ext == ".pdf":
        return extract_pdf_text(content, max_chars)
    elif ext == ".docx":
        return extract_docx_text(content, max_chars)
    elif ext == ".doc":
        # .doc은 바이너리 포맷 - docx로 변환 필요
        return "[.doc 파일은 지원되지 않습니다. .docx로 변환해주세요.]", False
    elif ext == ".xlsx":
        return extract_xlsx_text(content, max_chars)
    elif ext == ".xls":
        return extract_xls_text(content, max_chars)
    elif ext == ".pptx":
        return extract_pptx_text(content, max_chars)
    elif ext == ".ppt":
        return "[.ppt 파일은 지원되지 않습니다. .pptx로 변환해주세요.]", False
    elif ext == ".hwp":
        return extract_hwp_text(content, max_chars)
    elif ext == ".hwpx":
        return extract_hwpx_text(content, max_chars)
    elif ext == ".odt":
        return extract_odt_text(content, max_chars)
    elif ext in (".ods", ".odp"):
        # ODS/ODP도 odfpy로 처리 가능하지만 구조가 다름 - 향후 구현
        return extract_odt_text(content, max_chars)  # 기본 텍스트는 추출 가능
    elif ext == ".rtf":
        return extract_rtf_text(content, max_chars)
    else:
        return "", False

# 임시 파일 저장 경로
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

# 허용되는 파일 확장자
ALLOWED_EXTENSIONS = {
    # 텍스트 / 코드
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".css", ".js", ".py",
    ".java", ".cpp", ".c", ".h", ".ts", ".tsx", ".jsx", ".vue", ".svelte",
    ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala", ".r", ".m",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
    ".sql", ".graphql", ".yaml", ".yml", ".toml", ".ini", ".conf", ".cfg",
    ".env", ".gitignore", ".dockerfile", ".makefile",
    ".log", ".rst", ".tex", ".bib",
    # 문서
    ".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".ods", ".odp", ".rtf",
    ".hwp", ".hwpx",  # 한글 워드프로세서
    # 이미지
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg", ".bmp", ".ico", ".tiff", ".heic",
    # 오디오 (메타데이터만)
    ".mp3", ".wav", ".ogg", ".m4a", ".flac", ".aac",
    # 비디오 (메타데이터만)
    ".mp4", ".webm", ".avi", ".mov", ".mkv",
    # 아카이브 (메타데이터만)
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2",
    # 데이터
    ".parquet", ".pickle", ".pkl", ".npy", ".npz",
}

# 최대 파일 크기 (10MB)
MAX_FILE_SIZE = 10 * 1024 * 1024

# 메모리에 파일 컨텍스트 저장 (세션 기반)
uploaded_files: Dict[str, Dict[str, Any]] = {}


def get_file_type(filename: str) -> str:
    """파일 타입 분류"""
    ext = Path(filename).suffix.lower()

    IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg", ".bmp", ".ico", ".tiff", ".heic"}
    DOCUMENT_EXTENSIONS = {".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".odt", ".ods", ".odp", ".rtf", ".hwp", ".hwpx"}
    AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".m4a", ".flac", ".aac"}
    VIDEO_EXTENSIONS = {".mp4", ".webm", ".avi", ".mov", ".mkv"}
    ARCHIVE_EXTENSIONS = {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2"}
    DATA_EXTENSIONS = {".parquet", ".pickle", ".pkl", ".npy", ".npz"}

    if ext in IMAGE_EXTENSIONS:
        return "image"
    elif ext in DOCUMENT_EXTENSIONS:
        return "document"
    elif ext in AUDIO_EXTENSIONS:
        return "audio"
    elif ext in VIDEO_EXTENSIONS:
        return "video"
    elif ext in ARCHIVE_EXTENSIONS:
        return "archive"
    elif ext in DATA_EXTENSIONS:
        return "data"
    else:
        return "text"


@router.post("/files/upload")
async def upload_files(files: List[UploadFile] = File(...)):
    """파일 업로드 및 처리"""
    results = []

    for file in files:
        # 파일 확장자 검증
        ext = Path(file.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"지원하지 않는 파일 형식: {ext}"
            )

        # 파일 크기 검증
        content = await file.read()
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=400,
                detail=f"파일 크기가 너무 큽니다: {file.filename}"
            )

        # 파일 ID 생성
        file_id = str(uuid.uuid4())
        file_type = get_file_type(file.filename)

        # 파일 처리
        if file_type == "image":
            # 이미지는 base64로 인코딩
            file_data = {
                "id": file_id,
                "filename": file.filename,
                "type": file_type,
                "content_type": file.content_type,
                "base64": base64.b64encode(content).decode("utf-8"),
                "size": len(content)
            }
        elif file_type == "text":
            # 텍스트 파일은 내용 추출
            try:
                text_content = content.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    text_content = content.decode("cp949")  # Korean encoding
                except UnicodeDecodeError:
                    text_content = content.decode("latin-1")

            file_data = {
                "id": file_id,
                "filename": file.filename,
                "type": file_type,
                "content": text_content[:50000],  # 최대 50K 문자
                "size": len(content),
                "truncated": len(text_content) > 50000
            }
        elif file_type in ("audio", "video", "archive", "data"):
            # 미디어/아카이브/데이터 파일은 메타데이터만 저장
            file_path = UPLOAD_DIR / f"{file_id}{ext}"
            with open(file_path, "wb") as f:
                f.write(content)

            file_data = {
                "id": file_id,
                "filename": file.filename,
                "type": file_type,
                "path": str(file_path),
                "size": len(content),
                "content_type": file.content_type,
                "description": f"{file_type.capitalize()} file: {file.filename} ({len(content) / 1024:.1f} KB)"
            }
        else:
            # 문서 (PDF, DOCX, XLSX, HWP 등)는 파일 저장 + 텍스트 추출 시도
            file_path = UPLOAD_DIR / f"{file_id}{ext}"
            with open(file_path, "wb") as f:
                f.write(content)

            # 문서 텍스트 추출 (PDF, DOCX, XLSX, HWP 지원)
            extracted_text, truncated = extract_document_text(content, ext)

            file_data = {
                "id": file_id,
                "filename": file.filename,
                "type": file_type,
                "path": str(file_path),
                "size": len(content),
                "content": extracted_text if extracted_text else None,
                "truncated": truncated,
                "extraction_supported": bool(extracted_text and not extracted_text.startswith("["))
            }

        # 저장
        uploaded_files[file_id] = file_data

        # 결과 구성 (truncation 및 extraction 정보 포함)
        result_item = {
            "id": file_id,
            "filename": file.filename,
            "type": file_type,
            "size": len(content)
        }

        # 문서 파일의 경우 추출 상태 정보 추가
        if file_type == "document":
            result_item["truncated"] = file_data.get("truncated", False)
            result_item["extraction_supported"] = file_data.get("extraction_supported", False)
            if result_item["truncated"]:
                result_item["warning"] = "파일 내용이 너무 길어 일부만 추출되었습니다."
            elif not result_item["extraction_supported"]:
                result_item["warning"] = "이 파일 형식은 텍스트 추출이 지원되지 않습니다."

        # 텍스트 파일의 경우에도 truncation 정보 추가
        if file_type == "text" and file_data.get("truncated", False):
            result_item["truncated"] = True
            result_item["warning"] = "파일 내용이 너무 길어 일부만 추출되었습니다."

        results.append(result_item)

    return {"status": "ok", "files": results}


@router.get("/files/{file_id}")
async def get_file(file_id: str):
    """업로드된 파일 정보 조회"""
    if file_id not in uploaded_files:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다")

    file_data = uploaded_files[file_id]
    return file_data


@router.delete("/files/{file_id}")
async def delete_file(file_id: str):
    """업로드된 파일 삭제"""
    if file_id not in uploaded_files:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다")

    file_data = uploaded_files[file_id]

    # 파일 시스템에서 삭제
    if "path" in file_data:
        try:
            os.remove(file_data["path"])
        except FileNotFoundError:
            pass

    del uploaded_files[file_id]
    return {"status": "ok"}


def get_file_context(file_ids: List[str]) -> str:
    """파일들의 컨텍스트 문자열 생성 (AI 프롬프트용)"""
    contexts = []

    for file_id in file_ids:
        if file_id not in uploaded_files:
            continue

        file_data = uploaded_files[file_id]
        filename = file_data["filename"]
        file_type = file_data["type"]
        file_size = file_data.get("size", 0)
        size_str = f"{file_size / 1024:.1f} KB" if file_size < 1024 * 1024 else f"{file_size / (1024 * 1024):.1f} MB"

        if file_type == "text":
            content = file_data.get("content", "")
            truncated = file_data.get("truncated", False)
            context = f"[첨부파일: {filename}]\n{content}"
            if truncated:
                context += "\n...(내용이 잘림)"
            contexts.append(context)

        elif file_type == "image":
            # 이미지는 별도 처리 (get_file_images에서 Vision API로 전달)
            contexts.append(f"[첨부이미지: {filename} ({size_str})]")

        elif file_type == "document":
            # PDF 등 문서에서 추출된 텍스트가 있으면 포함
            content = file_data.get("content", "")
            if content:
                truncated = file_data.get("truncated", False)
                context = f"[첨부문서: {filename}]\n{content}"
                if truncated:
                    context += "\n...(내용이 잘림)"
                contexts.append(context)
            else:
                contexts.append(f"[첨부문서: {filename} ({size_str}) - 텍스트 추출 불가]")

        elif file_type == "audio":
            contexts.append(f"[첨부오디오: {filename} ({size_str})]")

        elif file_type == "video":
            contexts.append(f"[첨부비디오: {filename} ({size_str})]")

        elif file_type == "archive":
            contexts.append(f"[첨부압축파일: {filename} ({size_str})]")

        elif file_type == "data":
            contexts.append(f"[첨부데이터: {filename} ({size_str})]")

        else:
            contexts.append(f"[첨부파일: {filename} ({size_str})]")

    return "\n\n".join(contexts)


def get_file_images(file_ids: List[str]) -> List[Dict[str, Any]]:
    """이미지 파일들의 base64 데이터 반환 (Vision API용)"""
    images = []

    for file_id in file_ids:
        if file_id not in uploaded_files:
            continue

        file_data = uploaded_files[file_id]

        if file_data["type"] == "image":
            images.append({
                "filename": file_data["filename"],
                "base64": file_data.get("base64", ""),
                "content_type": file_data.get("content_type", "image/jpeg")
            })

    return images
